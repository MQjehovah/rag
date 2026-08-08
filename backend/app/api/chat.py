import json
import logging
import re
import urllib.parse
from typing import List, Dict, Any

import httpx
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.orm import Session

from app.config import settings
from app.core.rag import EmbeddingService, RerankerService
from app.core.retrieval import RetrievalPipeline
from app.models.database import Page
from app.api.deps import get_db
from app.core.jwt_utils import get_current_user

router = APIRouter(prefix="/api/chat", tags=["AI问答"])

logger = logging.getLogger(__name__)

_embedding_service = None
_reranker_service = None


def get_embedding_service():
    global _embedding_service
    if _embedding_service is None:
        _embedding_service = EmbeddingService()
    return _embedding_service


def get_reranker_service():
    global _reranker_service
    if _reranker_service is None:
        _reranker_service = RerankerService()
    return _reranker_service


class ChatRequest(BaseModel):
    query: str
    history: List[Dict[str, str]] = []


SYSTEM_PROMPT = """你是一个基于企业知识库的智能助手。请根据以下参考资料回答用户的问题。

要求：
- 只根据参考资料中的信息回答，不要编造内容
- 如果参考资料中没有相关信息，请诚实说明
- 回答时用 [1][2] 等编号标注对应资料，编号与参考资料列表一致
- 每个关键结论都要有编号引用，方便用户核对
- 如果参考资料附带的图片与回答直接相关，请在回答中直接嵌入图片，格式：![图片说明](图片URL)，只使用参考资料中出现的图片URL
- 使用中文回答"""

RAG_PROMPT_TEMPLATE = """参考资料：
{context}

---
{history}用户问题：{query}

请根据以上参考资料回答问题，并在相应位置用 [1][2] 等编号标注引用来源。"""


IMAGE_RE = re.compile(r'!\[[^\]]*\]\(([^)]+)\)')

SAFE_IMAGE_HOSTS = {"192.168.31.34", "192.168.31.8", "localhost", "127.0.0.1"}


def _normalize_image_url(u: str) -> str:
    """Rewrite external image URLs through our no-Referer proxy so images
    behind Referer checks (e.g. Alibaba OSS) still load in the browser."""
    if u.startswith("/"):
        return u
    if not (u.startswith("http://") or u.startswith("https://")):
        return u
    try:
        host = urllib.parse.urlsplit(u).hostname or ""
    except Exception:
        return u
    if host in SAFE_IMAGE_HOSTS:
        return u
    return "/api/upload/images/proxy?url=" + urllib.parse.quote(u, safe="")


def _extract_images(text: str, limit: int = 6) -> list:
    """Pull image URLs out of markdown content (skip huge base64 blobs)."""
    urls = []
    seen = set()
    for u in IMAGE_RE.findall(text or ""):
        u = u.strip()
        if not u or u in seen or u.startswith("data:"):
            continue
        norm = _normalize_image_url(u)
        if norm in seen:
            continue
        seen.add(norm)
        urls.append(norm)
        if len(urls) >= limit:
            break
    return urls


JUDGE_PROMPT = """判断已检索到的资料是否足够回答用户问题。

用户问题: {query}

已检索到的资料:
{notes}

只返回 JSON，不要其他内容:
{{
  "sufficient": true 或 false,
  "gap": "缺少的信息（一句话）",
  "followup_query": "还需要补充检索的子问题；如果已足够则为空字符串"
}}"""


async def _judge_sufficiency(query: str, notes: List[Dict[str, Any]]) -> Dict[str, Any]:
    parts = []
    for i, note in enumerate(notes, 1):
        chunks = note.get("chunks") or []
        excerpt = ""
        if chunks:
            excerpt = (chunks[0].get("content") or "")[:500]
        if not excerpt:
            excerpt = (note.get("content") or "")[:500]
        parts.append(f"[{i}] {note.get('title', '')}\n{excerpt}")
    return await _call_llm_json(
        [{"role": "user", "content": JUDGE_PROMPT.format(
            query=query,
            notes="\n\n".join(parts) if parts else "无",
        )}],
        context="agentic-judge",
    )


async def _agentic_search_notes(
    query: str,
    db: Session,
    current_user,
) -> List[Dict[str, Any]]:
    """Multi-hop retrieval: search, judge sufficiency, re-search if needed."""
    pipeline = RetrievalPipeline(
        db,
        embedding_svc=get_embedding_service(),
        reranker_svc=get_reranker_service(),
    )
    max_hops = max(settings.agentic_max_hops, 1)
    all_notes: List[Dict[str, Any]] = []
    seen_ids = set()
    current_q = query

    for hop in range(max_hops):
        outcome = await pipeline.retrieve(current_q, current_user, top_k=5)
        for note in outcome["results"]:
            if note["id"] not in seen_ids:
                seen_ids.add(note["id"])
                all_notes.append(note)
        # Close the read transaction before the judge LLM call / streaming, so
        # SQLite writers are never blocked by an idle read transaction.
        db.commit()

        if hop >= max_hops - 1 or not settings.llm_api_url:
            break
        try:
            decision = await _judge_sufficiency(query, all_notes)
        except Exception as e:
            logger.warning(f"Agentic judge failed: {e}")
            break
        followup = (decision.get("followup_query") or "").strip()
        if decision.get("sufficient") is True or not followup:
            break
        current_q = followup

    return all_notes


@router.post("")
async def chat(request: ChatRequest, db: Session = Depends(get_db), current_user=Depends(get_current_user)):
    if not settings.llm_api_url:
        raise HTTPException(status_code=500, detail="未配置 LLM API")

    notes = await _agentic_search_notes(request.query, db, current_user)

    context_parts = []
    sources = []
    for i, note in enumerate(notes, 1):
        chunks = note.get("chunks") or []
        chunk_texts = [c.get("content") or "" for c in chunks]
        excerpt = chunk_texts[0] if chunk_texts else (note.get("content") or "")[:500]
        if len(excerpt) > 2000:
            excerpt = excerpt[:2000] + "\n...(内容过长已截断)"
        imgs = _extract_images(" ".join(chunk_texts))
        if not imgs:
            imgs = _extract_images(note.get("content") or "")
        part = f"[{i}]《{note['title']}》\n{excerpt}"
        if imgs:
            part += f"\n相关图片: {', '.join(imgs[:4])}"
        context_parts.append(part)
        sources.append({
            "id": note["id"],
            "title": note["title"],
            "chunks": chunks[:3],
            "images": imgs[:6],
        })

    context = "\n\n".join(context_parts) if context_parts else "未找到相关笔记"
    history_msgs = []
    for h in (request.history or [])[-8:]:
        role = h.get("role")
        content = (h.get("content") or "").strip()
        if role in ("user", "assistant") and content:
            label = "用户" if role == "user" else "助手"
            history_msgs.append(f"{label}：{content[:1000]}")
    history_text = ("对话历史：\n" + "\n".join(history_msgs) + "\n\n") if history_msgs else ""
    user_message = RAG_PROMPT_TEMPLATE.format(
        context=context,
        query=request.query,
        history=history_text,
    )
    db.commit()  # no open read transaction while streaming the answer

    async def generate():
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                async with client.stream(
                    "POST",
                    settings.llm_api_url,
                    headers={
                        "Authorization": f"Bearer {settings.llm_api_key}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": settings.llm_model,
                        "messages": [
                            {"role": "system", "content": SYSTEM_PROMPT},
                            {"role": "user", "content": user_message},
                        ],
                        "stream": True,
                    },
                ) as resp:
                    resp.raise_for_status()
                    async for line in resp.aiter_lines():
                        if not line.startswith("data: "):
                            continue
                        data_str = line[6:]
                        if data_str.strip() == "[DONE]":
                            break
                        try:
                            data = json.loads(data_str)
                            delta = data.get("choices", [{}])[0].get("delta", {})
                            content = delta.get("content", "")
                            if content:
                                yield f"data: {json.dumps({'type': 'content', 'content': content}, ensure_ascii=False)}\n\n"
                        except json.JSONDecodeError:
                            continue

                yield f"data: {json.dumps({'type': 'sources', 'sources': sources}, ensure_ascii=False)}\n\n"
                yield "data: [DONE]\n\n"
        except httpx.HTTPStatusError as e:
            logger.error(f"LLM API error: {e}")
            yield f"data: {json.dumps({'type': 'error', 'content': f'LLM API 调用失败: {e.response.status_code}'}, ensure_ascii=False)}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as e:
            logger.error(f"Chat error: {e}")
            yield f"data: {json.dumps({'type': 'error', 'content': str(e)}, ensure_ascii=False)}\n\n"
            yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


class SaveNoteRequest(BaseModel):
    query: str
    answer: str


async def _call_llm_json(messages: list, context: str = "") -> dict:
    async with httpx.AsyncClient(timeout=120.0) as client:
        resp = await client.post(
            settings.llm_api_url,
            headers={
                "Authorization": f"Bearer {settings.llm_api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": settings.llm_model,
                "messages": messages,
                "stream": False,
            },
        )
        resp.raise_for_status()
        data = resp.json()
        content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
        logger.info(f"LLM raw response ({len(content)} chars): {content[:2000]}")
        if not content or not content.strip():
            logger.warning(f"LLM empty response [{context}]")
            return {}

        import re
        json_match = re.search(r'```(?:json)?\s*(\{[\s\S]*?\})\s*```', content)
        if json_match:
            content = json_match.group(1)

        try:
            return json.loads(content)
        except json.JSONDecodeError:
            pass

        brace_depth = 0
        start = -1
        for i, ch in enumerate(content):
            if ch == '{':
                if brace_depth == 0:
                    start = i
                brace_depth += 1
            elif ch == '}':
                brace_depth -= 1
                if brace_depth == 0 and start >= 0:
                    candidate = content[start:i + 1]
                    try:
                        result = json.loads(candidate)
                        logger.info(f"LLM parsed result [{context}]: {json.dumps(result, ensure_ascii=False)[:2000]}")
                        return result
                    except json.JSONDecodeError:
                        continue

        logger.warning(f"LLM response could not be parsed as JSON [{context}]: {content[:500]}")
        return {}


def _get_kb_context(db: Session) -> dict:
    notebooks = db.query(Notebook).all()
    nb_list = [{"id": nb.id, "name": nb.name} for nb in notebooks]
    pages = db.query(Page.id, Page.title, Page.notebook_id).order_by(Page.updated_at.desc()).limit(100).all()
    page_list = [{"id": p[0], "title": p[1], "notebook_id": p[2]} for p in pages]
    return {"notebooks": nb_list, "pages": page_list}


ORGANIZE_PROMPT = """你是一个企业知识库管理助手。请分析以下原始内容，完成知识整理和归类。

当前知识库结构：
笔记本列表：{kb_notebooks}
笔记列表（最近100篇）：{kb_pages}

原始内容（来源：{source}）：
{raw_content}

请完成以下任务：
1. 判断内容是否值得保存为知识笔记。如果内容毫无价值（如广告、导航栏、版权声明等），设置 should_save=false
2. 如果值得保存，将原始内容整理为结构化的 Markdown 格式，去除无效信息（广告、导航、版权、重复内容等），保留核心知识
3. 决定归类方式：
   - "create_notebook": 创建新笔记本 + 新笔记（当现有笔记本都不合适时）
   - "create_note": 在现有笔记本中创建新笔记
   - "update_note": 更新某篇已有笔记的内容（当内容是对已有笔记的补充时）

请以 JSON 格式返回：
{{
  "should_save": true/false,
  "action": "create_notebook" / "create_note" / "update_note",
  "title": "笔记标题",
  "notebook_id": "现有笔记本ID（create_note/update_note时）",
  "new_notebook_name": "新笔记本名称（create_notebook时）",
  "update_page_id": "要更新的笔记ID（update_note时）",
  "content": "整理后的Markdown内容",
  "summary": "100字以内摘要"
}}
只返回 JSON，不要其他内容。"""


@router.post("/save-note")
async def save_note(request: SaveNoteRequest, db: Session = Depends(get_db), current_user=Depends(get_current_user)):
    if not settings.llm_api_url:
        raise HTTPException(status_code=500, detail="未配置 LLM API")

    kb = _get_kb_context(db)
    prompt = ORGANIZE_PROMPT.format(
        kb_notebooks=json.dumps(kb["notebooks"], ensure_ascii=False),
        kb_pages=json.dumps(kb["pages"], ensure_ascii=False),
        source="AI对话",
        raw_content=f"用户问题：{request.query}\n\nAI回答：{request.answer}",
    )

    result = await _call_llm_json([{"role": "user", "content": prompt}], context="save-note")

    return {
        "should_save": result.get("should_save", True),
        "action": result.get("action", "create_note"),
        "title": result.get("title", request.query[:50]),
        "notebook_id": result.get("notebook_id"),
        "new_notebook_name": result.get("new_notebook_name"),
        "update_page_id": result.get("update_page_id"),
        "content": result.get("content", f"## {request.query}\n\n{request.answer}"),
        "summary": result.get("summary", ""),
        "notebooks": kb["notebooks"],
        "pages": kb["pages"],
    }


async def _extract_text_from_bytes(content: bytes, filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        import io
        try:
            import fitz
            import pymupdf4llm
            doc = fitz.open(stream=content, filetype="pdf")
            try:
                return pymupdf4llm.to_markdown(doc)
            finally:
                doc.close()
        except Exception:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
    elif ext == "docx":
        import io
        from docx import Document
        doc = Document(io.BytesIO(content))
        lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    elif ext in ("xlsx", "xlsm"):
        import io
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        lines = []
        try:
            for ws in wb.worksheets:
                if ws.max_row is None:
                    continue
                for row in ws.iter_rows(values_only=True):
                    vals = [str(v) for v in row if v is not None and str(v).strip()]
                    if vals:
                        lines.append(" | ".join(vals))
        finally:
            wb.close()
        return "\n".join(lines)
    elif ext == "pptx":
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            lines.append(t)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            lines.append(" | ".join(cells))
        return "\n".join(lines)
    elif ext in ("txt", "md", "csv", "json"):
        return content.decode("utf-8", errors="ignore")
    else:
        return content.decode("utf-8", errors="ignore")


def _looks_like_text(text: str) -> bool:
    """Reject binary/garbage content (failed file imports used to store raw
    file bytes as note content, which bloated the database and broke search)."""
    if not text or not text.strip():
        return False
    sample = text[:5000]
    if "\x00" in sample:
        return False
    control = sum(1 for ch in sample if ord(ch) < 32 and ch not in "\n\r\t")
    printable = sum(1 for ch in sample if ch.isprintable())
    if not sample:
        return False
    return control / len(sample) < 0.05 and printable / len(sample) > 0.6


def _extract_main_content(html: str) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "nav", "footer", "header", "aside",
                      "iframe", "noscript", "form", "button", "input", "select"]):
        tag.decompose()

    for tag_id in ["comments", "sidebar", "footer", "header", "nav", "menu",
                    "advertisement", "ad", "cookie", "banner", "popup", "modal"]:
        for el in soup.find_all(id=lambda x: x and tag_id in x.lower()):
            el.decompose()
    for cls in ["comment", "sidebar", "footer", "header", "nav", "menu",
                "ad", "advertisement", "cookie", "banner", "popup", "modal",
                "social", "share", "related", "recommend", "pagination"]:
        for el in soup.find_all(class_=lambda x: x and cls in " ".join(x).lower()):
            el.decompose()

    article = soup.find("article")
    if article:
        main = article
    else:
        candidates = soup.find_all(["div", "section", "main"])
        main = max(candidates, key=lambda el: len(el.get_text(strip=True))) if candidates else soup.body

    title = ""
    if soup.title and soup.title.string:
        title = soup.title.string.strip()

    for tag in main.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
        level = int(tag.name[1])
        tag.replace_with(f"\n{'#' * level} {tag.get_text(strip=True)}\n")

    for tag in main.find_all("p"):
        tag.replace_with(f"\n{tag.get_text(strip=True)}\n")

    for tag in main.find_all("li"):
        tag.replace_with(f"\n- {tag.get_text(strip=True)}")

    for tag in main.find_all("br"):
        tag.replace_with("\n")

    for tag in main.find_all(["a"]):
        href = tag.get("href", "")
        text = tag.get_text(strip=True)
        if href and text:
            tag.replace_with(f"[{text}]({href})")
        else:
            tag.replace_with(text)

    text = main.get_text()
    lines = [line.strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    while lines and not lines[0]:
        lines.pop(0)
    while lines and not lines[-1]:
        lines.pop()

    content = "\n".join(lines)
    return title, content


@router.post("/import/file")
async def import_file(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    current_user=Depends(get_current_user),
):
    if not settings.llm_api_url:
        raise HTTPException(status_code=500, detail="未配置 LLM API")

    content_bytes = await file.read()
    text = await _extract_text_from_bytes(content_bytes, file.filename or "file.txt")

    if not _looks_like_text(text):
        raise HTTPException(status_code=400, detail="无法从文件中提取有效文本（文件可能已损坏或格式不支持）")

    kb = _get_kb_context(db)
    prompt = ORGANIZE_PROMPT.format(
        kb_notebooks=json.dumps(kb["notebooks"], ensure_ascii=False),
        kb_pages=json.dumps(kb["pages"], ensure_ascii=False),
        source=f"文件: {file.filename}",
        raw_content=text[:6000],
    )

    result = await _call_llm_json([{"role": "user", "content": prompt}], context="import-file")

    return {
        "should_save": result.get("should_save", True),
        "action": result.get("action", "create_note"),
        "title": result.get("title", file.filename or "导入文件"),
        "notebook_id": result.get("notebook_id"),
        "new_notebook_name": result.get("new_notebook_name"),
        "update_page_id": result.get("update_page_id"),
        "content": result.get("content", text),
        "summary": result.get("summary", ""),
        "notebooks": kb["notebooks"],
        "pages": kb["pages"],
    }


class ImportUrlRequest(BaseModel):
    url: str


@router.post("/import/url")
async def import_url(request: ImportUrlRequest, db: Session = Depends(get_db), current_user=Depends(get_current_user)):
    if not settings.llm_api_url:
        raise HTTPException(status_code=500, detail="未配置 LLM API")

    try:
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            resp = await client.get(request.url)
            resp.raise_for_status()
            html = resp.text
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"抓取URL失败: {str(e)}")

    page_title, text = _extract_main_content(html)

    if not _looks_like_text(text):
        raise HTTPException(status_code=400, detail="无法提取有效的网页正文内容")

    kb = _get_kb_context(db)
    prompt = ORGANIZE_PROMPT.format(
        kb_notebooks=json.dumps(kb["notebooks"], ensure_ascii=False),
        kb_pages=json.dumps(kb["pages"], ensure_ascii=False),
        source=f"网页: {page_title or request.url}",
        raw_content=text[:6000],
    )

    result = await _call_llm_json([{"role": "user", "content": prompt}], context="import-url")

    return {
        "should_save": result.get("should_save", True),
        "action": result.get("action", "create_note"),
        "title": result.get("title", page_title or request.url),
        "notebook_id": result.get("notebook_id"),
        "new_notebook_name": result.get("new_notebook_name"),
        "update_page_id": result.get("update_page_id"),
        "content": result.get("content", text),
        "summary": result.get("summary", ""),
        "notebooks": kb["notebooks"],
        "pages": kb["pages"],
    }


class ImportTextRequest(BaseModel):
    text: str


@router.post("/import/text")
async def import_text(request: ImportTextRequest, db: Session = Depends(get_db), current_user=Depends(get_current_user)):
    if not settings.llm_api_url:
        raise HTTPException(status_code=500, detail="未配置 LLM API")

    if not request.text.strip():
        raise HTTPException(status_code=400, detail="文本内容为空")

    kb = _get_kb_context(db)
    prompt = ORGANIZE_PROMPT.format(
        kb_notebooks=json.dumps(kb["notebooks"], ensure_ascii=False),
        kb_pages=json.dumps(kb["pages"], ensure_ascii=False),
        source="用户粘贴文本",
        raw_content=request.text[:6000],
    )

    result = await _call_llm_json([{"role": "user", "content": prompt}], context="import-text")

    return {
        "should_save": result.get("should_save", True),
        "action": result.get("action", "create_note"),
        "title": result.get("title", "导入文本"),
        "notebook_id": result.get("notebook_id"),
        "new_notebook_name": result.get("new_notebook_name"),
        "update_page_id": result.get("update_page_id"),
        "content": result.get("content", request.text),
        "summary": result.get("summary", ""),
        "notebooks": kb["notebooks"],
        "pages": kb["pages"],
    }
